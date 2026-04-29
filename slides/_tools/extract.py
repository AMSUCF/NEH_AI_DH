"""Extract media, slide text, and speaker notes from .pptx files.

Outputs (per source deck, into slides/_extracted/<deck-stem>/):
  - media/        every embedded image and video, original filenames
  - slides.json   array of {index, title, body, notes, media[]}
  - inventory.md  human-readable summary used to build the week mapping

Usage (from repo root):
    python slides/_tools/extract.py slides/*.pptx
    python slides/_tools/extract.py slides/AI_Week_One.pptx slides/AI_WeekTwo.pptx

The .pptx itself is a zip; embedded media lives at ppt/media/. We use zipfile
for binary extraction and python-pptx to walk slides for titles/body/notes and
to map media to specific slide indices.
"""

from __future__ import annotations

import argparse
import json
import re
import sys
import zipfile
from dataclasses import dataclass, field, asdict
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu


REPO_ROOT = Path(__file__).resolve().parents[2]
DEFAULT_OUT = REPO_ROOT / "slides" / "_extracted"

IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".tiff", ".svg", ".emf", ".wmf"}
VIDEO_EXTS = {".mp4", ".mov", ".m4v", ".avi", ".wmv", ".webm", ".mkv"}
AUDIO_EXTS = {".mp3", ".wav", ".m4a", ".aac", ".ogg"}


@dataclass
class SlideRecord:
    index: int
    title: str = ""
    body: str = ""
    notes: str = ""
    media: list[str] = field(default_factory=list)


def classify(name: str) -> str:
    ext = Path(name).suffix.lower()
    if ext in IMAGE_EXTS:
        return "image"
    if ext in VIDEO_EXTS:
        return "video"
    if ext in AUDIO_EXTS:
        return "audio"
    return "other"


def human_size(num_bytes: int) -> str:
    for unit in ("B", "KB", "MB", "GB"):
        if num_bytes < 1024:
            return f"{num_bytes:.1f} {unit}"
        num_bytes /= 1024
    return f"{num_bytes:.1f} TB"


def extract_media_from_zip(pptx_path: Path, media_out: Path) -> dict[str, int]:
    """Copy ppt/media/* out of the .pptx zip. Returns {filename: byte_size}."""
    media_out.mkdir(parents=True, exist_ok=True)
    sizes: dict[str, int] = {}
    with zipfile.ZipFile(pptx_path) as z:
        for info in z.infolist():
            if not info.filename.startswith("ppt/media/"):
                continue
            name = Path(info.filename).name
            if not name:
                continue
            target = media_out / name
            with z.open(info) as src, open(target, "wb") as dst:
                dst.write(src.read())
            sizes[name] = info.file_size
    return sizes


def shape_text(shape) -> str:
    if not shape.has_text_frame:
        return ""
    parts = []
    for para in shape.text_frame.paragraphs:
        line = "".join(run.text for run in para.runs).strip()
        if line:
            parts.append(line)
    return "\n".join(parts)


def slide_title(slide) -> str:
    if slide.shapes.title and slide.shapes.title.has_text_frame:
        return slide.shapes.title.text.strip()
    # fallback: first non-empty text
    for shape in slide.shapes:
        text = shape_text(shape)
        if text:
            return text.split("\n", 1)[0].strip()[:120]
    return ""


def slide_body(slide) -> str:
    parts = []
    title = slide.shapes.title
    for shape in slide.shapes:
        if shape is title:
            continue
        text = shape_text(shape)
        if text:
            parts.append(text)
    return "\n\n".join(parts).strip()


def slide_notes(slide) -> str:
    if not slide.has_notes_slide:
        return ""
    nf = slide.notes_slide.notes_text_frame
    if nf is None:
        return ""
    return nf.text.strip()


MEDIA_REL_HINTS = ("/image", "/media", "/video", "/audio")


def collect_slide_media(slide) -> list[str]:
    """Return media filenames referenced from this slide, using relationships.

    Relying on `Picture.image.filename` is unreliable: python-pptx returns a
    generic "image.png" derived from the blob, not the unique `image7.png`
    that lives at `ppt/media/image7.png` in the zip. Instead we walk the
    slide part's relationships and resolve each media target by filename.
    """
    names: list[str] = []
    seen: set[str] = set()
    rels = slide.part.rels
    for rel in rels.values():
        reltype = rel.reltype.lower()
        if not any(h in reltype for h in MEDIA_REL_HINTS):
            continue
        target = rel.target_ref
        fname = Path(target).name
        if fname and fname not in seen:
            names.append(fname)
            seen.add(fname)
    return names


def extract_one(pptx_path: Path, out_root: Path) -> Path:
    deck_dir = out_root / pptx_path.stem
    media_dir = deck_dir / "media"
    deck_dir.mkdir(parents=True, exist_ok=True)

    print(f"[extract] {pptx_path.name}")
    sizes = extract_media_from_zip(pptx_path, media_dir)
    print(f"  media: {len(sizes)} files ({human_size(sum(sizes.values()))})")

    prs = Presentation(pptx_path)
    records: list[SlideRecord] = []
    for idx, slide in enumerate(prs.slides, start=1):
        rec = SlideRecord(
            index=idx,
            title=slide_title(slide),
            body=slide_body(slide),
            notes=slide_notes(slide),
            media=collect_slide_media(slide),
        )
        records.append(rec)
    print(f"  slides: {len(records)}")

    (deck_dir / "slides.json").write_text(
        json.dumps([asdict(r) for r in records], indent=2, ensure_ascii=False),
        encoding="utf-8",
    )

    write_inventory(deck_dir / "inventory.md", pptx_path.name, records, sizes)
    return deck_dir


def write_inventory(path: Path, source_name: str, records: list[SlideRecord], sizes: dict[str, int]) -> None:
    lines: list[str] = []
    lines.append(f"# Inventory: {source_name}")
    lines.append("")
    images = sum(1 for n in sizes if classify(n) == "image")
    videos = sum(1 for n in sizes if classify(n) == "video")
    audio = sum(1 for n in sizes if classify(n) == "audio")
    other = sum(1 for n in sizes if classify(n) == "other")
    lines.append(
        f"- Slides: **{len(records)}** | "
        f"Images: **{images}** | Videos: **{videos}** | Audio: **{audio}** | Other: **{other}**"
    )
    lines.append(f"- Media total size: **{human_size(sum(sizes.values()))}**")
    lines.append("")

    if videos:
        lines.append("## Videos (need YouTube/Vimeo URLs)")
        lines.append("")
        for name, size in sorted(sizes.items()):
            if classify(name) == "video":
                lines.append(f"- `{name}` ({human_size(size)})")
        lines.append("")

    lines.append("## Slide-by-slide")
    lines.append("")
    for rec in records:
        title = rec.title or "(untitled)"
        lines.append(f"### Slide {rec.index} — {title}")
        if rec.body:
            body = re.sub(r"\n+", " / ", rec.body).strip()
            if len(body) > 400:
                body = body[:400] + "…"
            lines.append(f"- Body: {body}")
        if rec.notes:
            notes = re.sub(r"\n+", " / ", rec.notes).strip()
            if len(notes) > 400:
                notes = notes[:400] + "…"
            lines.append(f"- Notes: {notes}")
        if rec.media:
            lines.append(f"- Media: {', '.join(f'`{m}`' for m in rec.media)}")
        lines.append("")
    path.write_text("\n".join(lines), encoding="utf-8")


def main(argv: list[str]) -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("pptx", nargs="+", type=Path, help="One or more .pptx paths")
    parser.add_argument("--out", type=Path, default=DEFAULT_OUT,
                        help=f"Output root (default: {DEFAULT_OUT})")
    args = parser.parse_args(argv)

    out_root: Path = args.out
    out_root.mkdir(parents=True, exist_ok=True)

    for p in args.pptx:
        if not p.exists():
            print(f"[skip] {p} (not found)", file=sys.stderr)
            continue
        try:
            extract_one(p, out_root)
        except Exception as e:
            print(f"[error] {p.name}: {e}", file=sys.stderr)
            raise
    print(f"\nDone. Output in {out_root}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))
